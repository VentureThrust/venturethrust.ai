# -*- coding: utf-8 -*-
"""
The sales deck for an angel network evaluating Deal Watch.

Built around the real workflow the investor already has: the deck arrives by
email, they read it, and one button on that page puts it on a watchlist. That
sequence is the whole pitch, because it answers the first objection any
network raises, which is that nobody wants to change how they work.

Screenshots live in scripts/bootcamp/shots/ and are picked up by filename.
Drop one in, re-run, and it lands on its slide. Anything missing renders as a
labelled placeholder so the deck is always presentable.

  01_email.png       the deck arriving in Gmail
  02_deck_watch.png  the deck open, with Add to Watchlist
  03_note.png        the note dialog
  04_watchlist.png   the watchlist
  05_report.png      a priority brief, open

  python scripts/bootcamp/make_sales_deck.py "Malabar Angel Network"
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "shots")
OUT = os.path.join(os.path.expanduser("~"), "Desktop", "om",
                   "VentureThrust_Sales_Deck.pptx")

# The live demo report. Put the real share link here and the deck carries a
# clickable "read the full report" line under the screenshot.
REPORT_LINK = os.environ.get("VT_REPORT_LINK", "")

NAVY = RGBColor(0x0D, 0x1B, 0x3E)
INK = RGBColor(0x0F, 0x17, 0x29)
CRIMSON = RGBColor(0x8B, 0x1E, 0x2D)
MUTED = RGBColor(0x6B, 0x72, 0x80)
RULE = RGBColor(0xD8, 0xDD, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF4, 0xF6, 0xF9)
GOLD = RGBColor(0xC7, 0xA2, 0x4B)

W, H = Inches(13.333), Inches(7.5)
L = Inches(0.9)
CW = Inches(11.53)


def box(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, first=False, after=0, before=0,
         line=None, align=PP_ALIGN.LEFT, italic=False, link=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    if link:
        r.hyperlink.address = link
    return p


def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def chrome(slide, n, kicker=None):
    rect(slide, L, Inches(6.92), CW, Pt(0.75), RULE)
    tf = box(slide, L, Inches(7.02), Inches(8), Inches(0.3))
    para(tf, "VentureThrust  ·  Deal Watch", 9, MUTED, first=True)
    tf2 = box(slide, Inches(11.6), Inches(7.02), Inches(0.85), Inches(0.3))
    para(tf2, str(n), 9, MUTED, first=True, align=PP_ALIGN.RIGHT)
    if kicker:
        t = box(slide, L, Inches(0.6), CW, Inches(0.3))
        para(t, kicker.upper(), 10.5, CRIMSON, bold=True, first=True)


def head(slide, title, y=Inches(0.98), size=30, sub=None):
    tf = box(slide, L, y, CW, Inches(0.9))
    para(tf, title, size, NAVY, bold=True, first=True, line=1.05)
    if sub:
        t2 = box(slide, L, y + Inches(0.62), CW, Inches(0.5))
        para(t2, sub, 15, MUTED, first=True, line=1.3)


def shot(slide, filename, x, y, w, h):
    """Letterbox a screenshot into the box, or draw a labelled placeholder."""
    path = os.path.join(SHOTS, filename)
    if os.path.exists(path):
        iw, ih = Image.open(path).size
        scale = min(w / iw, h / ih)
        dw, dh = int(iw * scale), int(ih * scale)
        dx, dy = x + int((w - dw) / 2), y + int((h - dh) / 2)
        slide.shapes.add_picture(path, dx, dy, dw, dh)
        rect(slide, dx, dy, dw, Pt(0.75), RULE)
        return True
    rect(slide, x, y, w, h, PALE, RULE)
    tf = box(slide, x, y + int(h / 2) - Inches(0.3), w, Inches(0.6))
    para(tf, "screenshot", 12, MUTED, first=True, align=PP_ALIGN.CENTER)
    para(tf, filename, 10, MUTED, align=PP_ALIGN.CENTER, italic=True)
    return False


def actor_tag(slide, actor, x=L, y=Inches(0.38)):
    """Who is doing this step. Never make the reader infer it from 'they'."""
    colour = {"FOUNDER": RGBColor(0x1E, 0x3A, 0x6E),
              "INVESTOR": CRIMSON,
              "VENTURETHRUST": RGBColor(0x2F, 0x6B, 0x4F)}.get(actor, NAVY)
    w = Inches(1.25 + 0.085 * len(actor))
    rect(slide, x, y, w, Inches(0.28), colour)
    tf = box(slide, x, y + Inches(0.045), w, Inches(0.24))
    para(tf, actor, 9.5, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
    return w


def step_slide(prs, n, step, title, sub, filename, actor="INVESTOR",
               note=None, link=None):
    """
    One workflow step. The screenshot is the slide, so the header is kept
    tight and the footer rule dropped: a shot letterboxed into a short box
    ends up small and marooned in white space.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])

    w = actor_tag(s, actor)
    t = box(s, L + w + Inches(0.22), Inches(0.4), CW, Inches(0.3))
    para(t, "STEP %d OF 5" % step, 10, MUTED, bold=True, first=True)

    # Step numeral in a filled square, left of the headline.
    rect(s, L, Inches(0.82), Inches(0.56), Inches(0.56), NAVY)
    tn = box(s, L, Inches(0.92), Inches(0.56), Inches(0.4))
    para(tn, str(step), 20, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)

    tf = box(s, Inches(1.62), Inches(0.8), Inches(10.8), Inches(0.55))
    para(tf, title, 26, NAVY, bold=True, first=True, line=1.05)
    t2 = box(s, Inches(1.62), Inches(1.33), Inches(10.8), Inches(0.4))
    para(t2, sub, 14, MUTED, first=True, line=1.25)

    top = Inches(1.85)
    bottom = Inches(7.22) if not (note or link) else Inches(6.82)
    shot(s, filename, L, top, CW, bottom - top)

    if note or link:
        tf3 = box(s, L, bottom + Inches(0.12), CW, Inches(0.4))
        if note:
            para(tf3, note, 13, INK, bold=True, first=True)
        if link:
            para(tf3, "Read the full sample report", 13, CRIMSON, bold=True,
                 first=not note, link=link)

    tp = box(s, Inches(12.3), Inches(7.05), Inches(0.6), Inches(0.3))
    para(tp, str(n), 9, MUTED, first=True, align=PP_ALIGN.RIGHT)
    return s


def bullets(slide, items, top=Inches(2.35), size=16.5, gap=13, width=CW):
    tf = box(slide, L, top, width, Inches(4.2))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            p.line_spacing = 1.25
            r = p.add_run(); r.text = lead + "  "
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = INK; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"
        else:
            para(tf, it, size, INK, first=(i == 0), after=gap, line=1.25)
    return tf


def build(network="your network"):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]
    n = 0

    # 1 ── Cover ───────────────────────────────────────────────────────────
    # The first line has to say what the buyer gets, not what is wrong with
    # their world. A B2B reader will not decode a clever hook, so the value
    # is the headline and the mechanism is demoted to the line under it.
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, Inches(0.42), H, NAVY)

    tf = box(s, Inches(1.35), Inches(1.55), Inches(10.6), Inches(3.2))
    para(tf, "DEAL WATCH", 12.5, CRIMSON, bold=True, first=True, after=18)
    # Two paragraphs rather than one with a line break: a manual break does
    # not take the paragraph line spacing, and the headline ends up airy.
    para(tf, "Get a second chance at the", 42, NAVY, bold=True, after=0, line=0.92)
    para(tf, "startups you passed on.", 42, NAVY, bold=True, after=22, line=0.92)
    para(tf, "We watch them for you and write to you only when the reason "
             "you said no is no longer true.", 17.5, MUTED, line=1.35)

    # Three outcomes on a hairline band, so the value is unmissable even if
    # nothing else on the page is read.
    rect(s, Inches(1.35), Inches(4.62), Inches(10.6), Pt(0.75), RULE)
    outcomes = [
        ("Deal flow you already own", "No new sourcing. These are companies you have met."),
        ("One click to start", "Nothing about how your members work changes."),
        ("Silence by default", "Most months you hear nothing, and that is the point."),
    ]
    colw = Inches(3.53)
    for i, (a, b) in enumerate(outcomes):
        x = Inches(1.35) + i * colw
        if i:
            rect(s, x - Inches(0.16), Inches(4.85), Pt(0.75), Inches(0.86), RULE)
        t = box(s, x, Inches(4.88), colw - Inches(0.35), Inches(1.0))
        para(t, a, 13.5, NAVY, bold=True, first=True, after=5)
        para(t, b, 11.5, MUTED, line=1.3)

    rect(s, Inches(1.35), Inches(6.08), Inches(10.6), Pt(0.75), RULE)
    tf2 = box(s, Inches(1.35), Inches(6.32), Inches(10.6), Inches(0.9))
    para(tf2, "Prepared for %s" % network, 14.5, INK, bold=True, first=True, after=4)
    para(tf2, "Omprakash Borkar  ·  VentureThrust  ·  venturethrust.com", 12.5, MUTED)

    # 2 ── The gap ─────────────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank); chrome(s, n, "The gap")
    head(s, "An investor says no once, and then goes blind")
    bullets(s, [
        ("The no is usually right.", "Too early, wrong economics, no clearance yet."),
        ("The startup keeps building anyway.", "Six months later the reason for the no is gone."),
        ("Nobody tells the investor.", "They find out from a funding announcement, at a higher price."),
        ("So the second look never happens.", "Not because the deal got worse, but because nobody was watching."),
    ], top=Inches(2.25))
    rect(s, L, Inches(5.6), CW, Pt(0.75), RULE)
    tf = box(s, L, Inches(5.85), CW, Inches(0.9))
    para(tf, "A network of a hundred members passes on hundreds of startups a year. "
             "Not one of them is being watched.", 17, NAVY, bold=True, first=True, line=1.3)

    # 3 ── Nothing changes ─────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank); chrome(s, n, "The important part")
    head(s, "Nothing about how your investors work changes")
    bullets(s, [
        ("The founder still emails the deck.", "Exactly as they do today."),
        ("The investor still opens it and reads it.", "Same inbox, same habit."),
        ("One extra button on that page.", "That is the entire change for the investor."),
        ("VentureThrust does the rest.", "No new login, no forms, no process to adopt."),
    ], top=Inches(2.3))
    rect(s, L, Inches(5.5), CW, Pt(0.75), RULE)
    tf = box(s, L, Inches(5.78), CW, Inches(0.9))
    para(tf, "Every tool that asks an investor to change their workflow dies in the first "
             "week. This one asks the investor for a single click.",
         17, NAVY, bold=True, first=True, line=1.3)

    # 4 to 8 ── The workflow ───────────────────────────────────────────────
    n += 1
    step_slide(prs, n, 1, "The founder sends the deck by email",
               "Straight to the investor's inbox, exactly as it happens today.",
               "01_email.png", actor="FOUNDER")

    n += 1
    step_slide(prs, n, 2, "The investor opens it and sees one button",
               "Add to Watchlist sits on the deck the investor is already reading.",
               "02_deck_watch.png", actor="INVESTOR")

    n += 1
    step_slide(prs, n, 3, "The investor writes why they passed",
               "Optional. This note is what the brief will be measured against later.",
               "03_note.png", actor="INVESTOR",
               note="Quarterly reports are opt in, per startup. Off by default.")

    n += 1
    step_slide(prs, n, 4, "The startup sits on the investor's watchlist",
               "The investor does nothing further. No forms, no reminders, no chasing.",
               "04_watchlist.png", actor="VENTURETHRUST")

    # Step 5 gets its own layout. The brief is the deliverable, so the page
    # runs full height on the left with what is in it spelled out beside it.
    n += 1
    s = prs.slides.add_slide(blank)
    w = actor_tag(s, "VENTURETHRUST")
    t = box(s, L + w + Inches(0.22), Inches(0.4), CW, Inches(0.3))
    para(t, "STEP 5 OF 5", 10, MUTED, bold=True, first=True)
    rect(s, L, Inches(0.82), Inches(0.56), Inches(0.56), NAVY)
    tn = box(s, L, Inches(0.92), Inches(0.56), Inches(0.4))
    para(tn, "5", 20, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
    tf = box(s, Inches(1.62), Inches(0.8), Inches(10.8), Inches(0.55))
    para(tf, "We send the investor a brief, and only when it matters",
         26, NAVY, bold=True, first=True, line=1.05)
    t2 = box(s, Inches(1.62), Inches(1.33), Inches(10.8), Inches(0.4))
    para(t2, "Not on a schedule. Only when the reason the investor passed stops being true.",
         14, MUTED, first=True, line=1.25)

    shot(s, "05_report_page.png", L, Inches(1.9), Inches(4.6), Inches(5.2))

    tx = Inches(6.0)
    tw = Inches(6.6)
    tf = box(s, tx, Inches(2.0), tw, Inches(4.2))
    for i, (a, b) in enumerate([
        ("Why you are getting this",
         "Quoted from the note the investor wrote when they passed."),
        ("Then versus now",
         "The same metrics, at the pass and today, with the arithmetic shown."),
        ("What changed and what did not",
         "Durable drivers separated from seasonal luck."),
        ("What we would still watch",
         "The thing that could still go wrong."),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3); p.line_spacing = 1.2
        r = p.add_run(); r.text = a
        r.font.size = Pt(15.5); r.font.bold = True; r.font.color.rgb = INK
        r.font.name = "Calibri"
        p2 = tf.add_paragraph(); p2.space_after = Pt(15); p2.line_spacing = 1.25
        r2 = p2.add_run(); r2.text = b
        r2.font.size = Pt(13.5); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

    rect(s, tx, Inches(5.55), tw, Pt(0.75), RULE)
    tf2 = box(s, tx, Inches(5.8), tw, Inches(1.0))
    para(tf2, "Every brief ends: we explain, the investor decides.", 15.5, CRIMSON,
         bold=True, first=True, after=8)
    if REPORT_LINK:
        para(tf2, "Read the full sample report", 14, NAVY, bold=True,
             link=REPORT_LINK)
    tp = box(s, Inches(12.3), Inches(7.05), Inches(0.6), Inches(0.3))
    para(tp, str(n), 9, MUTED, first=True, align=PP_ALIGN.RIGHT)

    # 9 ── The three layers ────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank); chrome(s, n, "What happens in between")
    head(s, "Software catches it. AI reads it. A human confirms it.")
    tf = box(s, L, Inches(2.15), CW, Inches(2.2))
    for i, (a, b) in enumerate([
        ("1. The software catches it", "A document changes in the founder's data room."),
        ("2. The AI reads it", "It compares the change against the investor's own note."),
        ("3. A human confirms it", "Nothing reaches the investor until a person agrees it matters."),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = a + "   "
        r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = CRIMSON
        r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = b
        r2.font.size = Pt(17); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"
    rect(s, L, Inches(4.35), CW, Pt(0.75), RULE)
    bullets(s, [
        ("The founder consents.", "The startup keeps its room on VentureThrust and controls what is shared. Nothing is scraped."),
        ("The third layer is the product.", "One useless alert and the investor stops reading the next one. So a person signs off on every brief."),
    ], top=Inches(4.7), size=15.5, gap=12)

    # 10 ── Silence ────────────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank); chrome(s, n, "Why investors trust it")
    head(s, "Silence is the product")
    bullets(s, [
        "Nothing happened this month, so the investor hears nothing. That is the default.",
        "A quarterly report goes only to the startups where the investor asked for one.",
        "Most months an investor hears from us once, or not at all.",
        "Anyone can send more email. The scarce thing is a service that stays quiet.",
    ], top=Inches(2.3))
    rect(s, L, Inches(5.5), CW, Pt(0.75), RULE)
    tf = box(s, L, Inches(5.78), CW, Inches(0.9))
    para(tf, "Every brief ends the same way: we explain, the investor decides. "
             "We never say invest.", 17, NAVY, bold=True, first=True, line=1.3)

    # 11 ── Pricing ────────────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank); chrome(s, n, "Commercials")
    head(s, "What it costs")
    cols = [("$149", "per member, per month"),
            ("Rs 1.5 lakh", "per member, per year"),
            ("One", "recovered deal pays for a decade")]
    colw = int(CW / 3)
    for i, (v, lab) in enumerate(cols):
        x = L + i * colw
        if i:
            rect(s, x - Inches(0.2), Inches(2.3), Pt(0.75), Inches(1.5), RULE)
        tf = box(s, x, Inches(2.3), Emu(colw - Inches(0.4)), Inches(1.6))
        para(tf, v, 38, NAVY, bold=True, first=True, after=6)
        para(tf, lab, 13, MUTED, line=1.25)
    rect(s, L, Inches(4.25), CW, Pt(0.75), RULE)
    bullets(s, [
        ("Network pricing.", "Volume terms where a network takes seats for its members."),
        ("No setup fee, no lock in.", "Monthly, cancel whenever."),
        ("The economics are not subtle.", "One deal re-entered at the earlier price covers many years of this."),
    ], top=Inches(4.55), size=15.5, gap=11)

    # 12 ── The offer ──────────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank); chrome(s, n, "How to start")
    head(s, "Three names, no money")
    bullets(s, [
        ("Send me three startups.", "Ones you or your members passed on in the last year."),
        ("I will send briefs on all three, free.", "Written exactly as you saw, inside one week."),
        ("Show them to your members.", "If nobody finds them useful, you have lost nothing."),
        ("Then we talk about seats.", "Not before."),
    ], top=Inches(2.3))
    rect(s, L, Inches(5.5), CW, Pt(0.75), RULE)
    tf = box(s, L, Inches(5.78), CW, Inches(0.9))
    para(tf, "This is a pilot, not a purchase. Nothing to approve and nothing to sign.",
         17, NAVY, bold=True, first=True)

    # 13 ── Close ──────────────────────────────────────────────────────────
    n += 1; s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, NAVY)
    tf = box(s, Inches(1.5), Inches(2.5), Inches(10.4), Inches(2.8))
    para(tf, "We explain. You decide.", 44, WHITE, bold=True, first=True, after=26)
    para(tf, "Omprakash Borkar", 17, GOLD, bold=True, after=6)
    para(tf, "omprakash@venturethrust.com", 15, RGBColor(0x9D, 0xB2, 0xD9), after=4)
    para(tf, "venturethrust.com  ·  venturethrust.com/investors", 15,
         RGBColor(0x9D, 0xB2, 0xD9))

    os.makedirs(SHOTS, exist_ok=True)
    prs.save(OUT)

    needed = ["01_email.png", "02_deck_watch.png", "03_note.png",
              "04_watchlist.png", "05_report.png"]
    missing = [f for f in needed if not os.path.exists(os.path.join(SHOTS, f))]
    print("DECK:  %s" % OUT)
    print("SHOTS: %s" % SHOTS)
    print("have:  %d of %d" % (len(needed) - len(missing), len(needed)))
    if missing:
        print("need:  %s" % ", ".join(missing))
    if not REPORT_LINK:
        print("note:  set VT_REPORT_LINK to add the clickable full report link")
    return OUT


if __name__ == "__main__":
    build(sys.argv[1] if len(sys.argv) > 1 else "your network")
