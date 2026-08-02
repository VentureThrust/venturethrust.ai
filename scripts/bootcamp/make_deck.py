# -*- coding: utf-8 -*-
"""
The deck for the IIMK LIVE IDEA VAULT bootcamp one on one mentoring sessions,
3 August 2026.

Built for a 20 minute conversation, not a keynote. Ten slides, five minutes
of talking, fifteen minutes of them talking back. That ratio is the point.

  python scripts/bootcamp/make_deck.py
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.join(os.path.expanduser("~"), "Desktop", "om",
                   "VentureThrust_Bootcamp_Deck.pptx")

NAVY = RGBColor(0x0D, 0x1B, 0x3E)
CRIMSON = RGBColor(0x8B, 0x1E, 0x2D)
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
RULE = RGBColor(0xD1, 0xD5, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF4, 0xB4, 0x00)

W, H = Inches(13.333), Inches(7.5)
L = Inches(0.95)          # left margin
CW = Inches(11.4)         # content width


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, space_after=0, space_before=0,
         first=False, align=PP_ALIGN.LEFT, italic=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def rule(slide, y, x=L, w=CW, color=RULE, thick=Pt(1)):
    from pptx.enum.shapes import MSO_SHAPE
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def footer(slide, n):
    tf = box(slide, L, Inches(6.85), CW, Inches(0.35))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "VentureThrust  ·  Omprakash Borkar  ·  IIMK LIVE IDEA VAULT, 3 August 2026"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED
    r.font.name = "Calibri"
    tf2 = box(slide, Inches(11.6), Inches(6.85), Inches(0.8), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(n)
    r2.font.size = Pt(9)
    r2.font.color.rgb = MUTED
    r2.font.name = "Calibri"


def title_slide(prs, kicker, title, subs):
    s = blank(prs)
    from pptx.enum.shapes import MSO_SHAPE
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.42), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = box(s, Inches(1.4), Inches(2.35), Inches(10.6), Inches(3.2))
    para(tf, kicker, 13, CRIMSON, bold=True, first=True, space_after=14)
    para(tf, title, 46, NAVY, bold=True, space_after=16, line=1.05)
    for sline in subs:
        para(tf, sline, 16, MUTED, space_after=6)
    return s


def content_slide(prs, n, head, kicker=None):
    s = blank(prs)
    y = Inches(0.7)
    if kicker:
        tf = box(s, L, y, CW, Inches(0.32))
        para(tf, kicker.upper(), 11, CRIMSON, bold=True, first=True)
        y = Inches(1.06)
    tf = box(s, L, y, CW, Inches(0.95))
    para(tf, head, 32, NAVY, bold=True, first=True, line=1.06)
    rule(s, Inches(2.12))
    footer(s, n)
    return s


def bullets(slide, items, top=Inches(2.5), size=17, gap=15):
    tf = box(slide, L, top, CW, Inches(4.0))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            p.line_spacing = 1.25
            r = p.add_run(); r.text = lead + "  "
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = INK; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED
            r2.font.name = "Calibri"
        else:
            para(tf, it, size, INK, first=(i == 0), space_after=gap, line=1.25)
    return tf


def stat_row(slide, stats, top=Inches(2.6)):
    """Big numbers across the slide, hairline separated."""
    n = len(stats)
    colw = int(CW / n)
    for i, (value, label) in enumerate(stats):
        x = L + i * colw
        tf = box(slide, x, top, Emu(colw - Inches(0.3)), Inches(1.6))
        para(tf, value, 40, NAVY, bold=True, first=True, space_after=4)
        para(tf, label, 12.5, MUTED, line=1.2)
        if i:
            rule(slide, top, x=x - Inches(0.18), w=Pt(1), color=RULE, thick=Inches(1.5))


def quote_slide(prs, n, text, attrib):
    s = blank(prs)
    from pptx.enum.shapes import MSO_SHAPE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    tf = box(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(3.0))
    para(tf, text, 34, WHITE, bold=True, first=True, line=1.18, space_after=22)
    para(tf, attrib, 15, RGBColor(0x9D, 0xB2, 0xD9))
    return s


def build():
    prs = deck()

    # 1 ────────────────────────────────────────────────────────────────────
    title_slide(
        prs,
        "IIMK LIVE  ·  IDEA VAULT BOOTCAMP  ·  3 AUGUST 2026",
        "VentureThrust",
        ["The investor's second look.",
         "Omprakash Borkar  ·  Founder  ·  Nagpur",
         "omprakash@venturethrust.com  ·  venturethrust.com"],
    )

    # 2 ────────────────────────────────────────────────────────────────────
    quote_slide(
        prs, 2,
        "You have seen thousands of pitches.\nWhat happened to the ones you said no to?",
        "This is the only question my product answers.",
    )

    # 3 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 3, "An investor says no once, and then goes blind",
                      kicker="The gap")
    bullets(s, [
        ("The no is usually right.", "Too early, wrong economics, no clearance yet."),
        ("The startup keeps building anyway.", "Six months later the reason for the no is gone."),
        ("Nobody tells the investor.", "They find out from a funding announcement, at a higher price."),
        ("So the second look never happens.", "Not because the deal got worse. Because nobody was watching."),
    ])

    # 4 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 4, "Deal Watch", kicker="What I built")
    tf = box(s, L, Inches(2.45), CW, Inches(0.9))
    para(tf, "An investor puts the startups they passed on to a watchlist, "
             "with the reason they passed in their own words.",
         20, INK, first=True, line=1.3)
    para(tf, "We tell them the day that reason stops being true.", 20, INK,
         bold=True, line=1.3, space_before=10)

    rule(s, Inches(4.15))
    layers = [
        ("1. The software catches it", "A document changes in the startup's data room."),
        ("2. The AI reads it", "It compares what changed against that investor's own note."),
        ("3. A human confirms it", "Nothing reaches the investor until a person agrees it matters."),
    ]
    tf = box(s, L, Inches(4.45), CW, Inches(2.0))
    for i, (a, b) in enumerate(layers):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        r = p.add_run(); r.text = a + "   "
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = CRIMSON
        r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = b
        r2.font.size = Pt(16); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

    # 5 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 5, "What actually lands in their inbox",
                      kicker="The product is the brief")
    bullets(s, [
        ("Why you are getting this.", "Quoted from the note they wrote when they passed."),
        ("Then versus now.", "The same metrics, at the pass and today, with the arithmetic shown."),
        ("What changed and what did not.", "Durable drivers separated from seasonal luck."),
        ("What we would still watch.", "The thing that could still go wrong."),
    ], top=Inches(2.45))
    tf = box(s, L, Inches(5.55), CW, Inches(0.9))
    para(tf, "Every brief ends the same way: we explain, you decide.", 19,
         NAVY, bold=True, first=True)
    para(tf, "We never say invest. The moment we do, we are a tipster and not a service.",
         14, MUTED, space_before=6)

    # 6 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 6, "Silence is the product", kicker="Why they trust it")
    bullets(s, [
        "Nothing happened this month, so we send nothing. That is the default.",
        "A quarterly report goes out only on the startups where the investor asked for one.",
        "Most months an investor hears from us once, or not at all.",
        "Anyone can send more email. The scarce thing is a service that stays quiet.",
    ], top=Inches(2.5))

    # 7 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 7, "The arithmetic an angel already knows",
                      kicker="Why it is worth paying for")
    stat_row(s, [
        ("$149", "per month, per investor"),
        ("Rs 1.5 L", "a year, all in"),
        ("One", "deal recovered pays for a decade"),
    ], top=Inches(2.6))
    tf = box(s, L, Inches(4.6), CW, Inches(1.6))
    para(tf, "An angel network of a hundred members passes on hundreds of startups a year.",
         18, INK, first=True, space_after=10, line=1.3)
    para(tf, "Nobody is watching any of them. If one comes back at the right moment, "
             "the whole network shares the return.", 18, INK, line=1.3)

    # 8 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 8, "Where I actually am", kicker="No decoration")
    bullets(s, [
        ("Live product.", "venturethrust.com. Data room, sharing, analytics, Deal Watch, all working today."),
        ("Zero paying investors.", "That is the honest number, and it is why I am in this room."),
        ("Built solo.", "Sole proprietorship in Nagpur. No team to pay, so the runway is long."),
        ("One investor conversation so far.", "It was with IIMK LIVE, and it changed the product."),
    ], top=Inches(2.5))

    # 9 ────────────────────────────────────────────────────────────────────
    s = content_slide(prs, 9, "What I am asking you for", kicker="The ask")
    bullets(s, [
        ("Not a subscription.", "I am not trying to sell you a seat in this meeting."),
        ("Three names.", "Three startups you or your network passed on in the last year."),
        ("I will send briefs on all three, free.", "Written the way you just saw. Inside one week."),
        ("Then you decide.", "Show them to your members. If nobody finds it useful, you have lost nothing."),
    ], top=Inches(2.5))

    # 10 ───────────────────────────────────────────────────────────────────
    s = blank(prs)
    from pptx.enum.shapes import MSO_SHAPE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background(); bg.shadow.inherit = False
    tf = box(s, Inches(1.5), Inches(2.6), Inches(10.3), Inches(2.6))
    para(tf, "We explain. You decide.", 44, WHITE, bold=True, first=True, space_after=24)
    para(tf, "Omprakash Borkar", 18, GOLD, bold=True, space_after=6)
    para(tf, "omprakash@venturethrust.com", 16, RGBColor(0x9D, 0xB2, 0xD9), space_after=4)
    para(tf, "venturethrust.com  ·  venturethrust.com/investors", 16,
         RGBColor(0x9D, 0xB2, 0xD9))

    prs.save(OUT)
    print("DECK: %s" % OUT)
    return OUT


if __name__ == "__main__":
    build()
